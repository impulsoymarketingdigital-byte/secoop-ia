#!/usr/bin/env python3
"""
JIREHAI - SECOP AI Analyst
parser.py — Extrae texto y tablas de documentos SECOP II

Soporta: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, imágenes (via OCR)

Estrategias por tipo:
  PDF   → PyMuPDF (texto) + pdfplumber (tablas precisas)
  DOCX  → python-docx (párrafos + tablas con celdas fusionadas)
  XLSX  → openpyxl (respeta fusiones) + pandas (fallback)
  PPTX  → python-pptx (texto de cada diapositiva + tablas)
  IMG   → marcado como escaneado → OCR en fase siguiente
"""
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

SCANNED_PAGE_THRESHOLD = 50   # chars/página mínimos para no ser "escaneado"


# ─────────────────────────────────────────────────────────────────────────────
# Punto de entrada
# ─────────────────────────────────────────────────────────────────────────────

def parse_document(file_path) -> dict:
    path = Path(file_path)
    if not path.exists():
        return _empty(error=f"Archivo no encontrado: {file_path}")

    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _parse_pdf(path)
        elif suffix in (".docx", ".doc"):
            return _parse_docx(path)
        elif suffix in (".xlsx", ".xls"):
            return _parse_excel(path)
        elif suffix in (".pptx", ".ppt"):
            return _parse_pptx(path)
        elif suffix in (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"):
            return _parse_image(path)
        else:
            # Intentar como texto plano
            try:
                text = path.read_text(encoding="utf-8", errors="replace")
                return _empty(file_type=suffix) | {
                    "text": text, "page_count": 1, "is_scanned": False,
                    "pages": [{"page": 1, "text": text, "char_count": len(text)}],
                    "error": None,
                }
            except Exception:
                return _empty(error=f"Tipo no soportado: {suffix}", file_type=suffix)
    except Exception as e:
        logger.exception(f"Error parseando {path.name}: {e}")
        return _empty(error=str(e), file_type=suffix)


# ─────────────────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> dict:
    result = _empty(file_type=".pdf")
    pages_data = []
    all_text_parts = []

    # ── PyMuPDF: texto página a página ──────────────────────────────────
    try:
        import fitz
        doc = fitz.open(str(path))
        result["page_count"] = len(doc)

        for pn in range(len(doc)):
            page = doc[pn]
            # Extraer con preservación de layout para tablas simples
            text = page.get_text("text") or ""
            text = _clean_text(text)
            pages_data.append({"page": pn + 1, "text": text, "char_count": len(text.strip())})
            if text.strip():
                all_text_parts.append(f"[Página {pn+1}]\n{text}")
        doc.close()
    except ImportError:
        logger.warning("PyMuPDF no disponible")
    except Exception as e:
        logger.warning(f"PyMuPDF error: {e}")

    # ── Detectar si es escaneado ─────────────────────────────────────────
    if pages_data:
        nonempty = [p for p in pages_data if p["char_count"] > 0]
        avg = sum(p["char_count"] for p in nonempty) / max(len(nonempty), 1)
        result["is_scanned"] = avg < SCANNED_PAGE_THRESHOLD
    else:
        result["is_scanned"] = True

    # ── pdfplumber: tablas + complementar texto vacío ────────────────────
    tables_extracted = []
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            if not result["page_count"]:
                result["page_count"] = len(pdf.pages)

            for pn, page in enumerate(pdf.pages):
                # Tablas
                try:
                    for t_idx, table in enumerate(page.extract_tables() or []):
                        if not table or len(table) < 2:
                            continue
                        # Buscar fila de encabezado real (primera fila no vacía)
                        header_row_idx = 0
                        for ri, row in enumerate(table):
                            if any(c and str(c).strip() for c in row):
                                header_row_idx = ri
                                break
                        headers = [str(c or "").strip() for c in table[header_row_idx]]
                        rows = []
                        for row in table[header_row_idx + 1:]:
                            cleaned = [str(c or "").strip() for c in row]
                            if any(c for c in cleaned):
                                rows.append(cleaned)
                        if rows:
                            tables_extracted.append({
                                "source": f"pdf_p{pn+1}_t{t_idx+1}",
                                "headers": headers,
                                "data": rows,
                                "page": pn + 1,
                            })
                except Exception as te:
                    logger.debug(f"Tabla PDF p{pn+1}: {te}")

                # Texto faltante
                if pn < len(pages_data) and pages_data[pn]["char_count"] < 10:
                    try:
                        pt = _clean_text(page.extract_text() or "")
                        if pt.strip():
                            pages_data[pn]["text"] = pt
                            pages_data[pn]["char_count"] = len(pt.strip())
                            all_text_parts.append(f"[Página {pn+1}]\n{pt}")
                    except Exception:
                        pass
    except ImportError:
        logger.warning("pdfplumber no disponible")
    except Exception as e:
        logger.warning(f"pdfplumber error: {e}")

    result.update({
        "pages": pages_data,
        "tables": tables_extracted,
        "text": "\n\n".join(all_text_parts),
        "error": None,
    })
    return result


# ─────────────────────────────────────────────────────────────────────────────
# Word (DOCX / DOC)
# ─────────────────────────────────────────────────────────────────────────────

def _parse_docx(path: Path) -> dict:
    result = _empty(file_type=path.suffix.lower())
    try:
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(str(path))
        parts = []

        # Párrafos (incluyendo encabezados con estilo)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Detectar si es título/sección
                style_name = (para.style.name or "").lower()
                if "heading" in style_name or "título" in style_name:
                    parts.append(f"\n## {text}")
                else:
                    parts.append(text)

        # Tablas (con manejo de celdas fusionadas)
        tables_extracted = []
        for t_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.strip().replace("\n", " "))
                if any(c for c in cells):
                    rows.append(cells)

            # Deduplicar filas repetidas (efecto de fusión)
            unique_rows = []
            prev = None
            for r in rows:
                if r != prev:
                    unique_rows.append(r)
                    prev = r

            if unique_rows:
                headers = unique_rows[0] if unique_rows else []
                data = unique_rows[1:] if len(unique_rows) > 1 else []
                # Agregar texto de tabla al texto principal
                table_text = _table_to_text(headers, data, f"Tabla {t_idx+1}")
                parts.append(table_text)
                tables_extracted.append({
                    "source": f"docx_t{t_idx+1}",
                    "headers": headers,
                    "data": data,
                    "page": None,
                })

        full_text = "\n".join(parts)
        result.update({
            "text": full_text,
            "pages": [{"page": 1, "text": full_text, "char_count": len(full_text)}],
            "tables": tables_extracted,
            "page_count": 1,
            "is_scanned": False,
            "error": None,
        })
    except ImportError:
        result["error"] = "python-docx no disponible — instalar: pip install python-docx"
    except Exception as e:
        result["error"] = str(e)
        logger.warning(f"DOCX parse error {path.name}: {e}")
    return result


# ─────────────────────────────────────────────────────────────────────────────
# Excel (XLSX / XLS) — manejo robusto de layouts irregulares
# ─────────────────────────────────────────────────────────────────────────────

def _parse_excel(path: Path) -> dict:
    result = _empty(file_type=path.suffix.lower())
    tables_extracted = []
    text_parts = []

    # ── openpyxl: lectura fiel incluyendo celdas fusionadas ─────────────
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            if ws.max_row == 0 or ws.max_column == 0:
                continue

            # Desenfusionar: crear grilla plana con valores propagados
            grid = _unfuse_sheet(ws)
            if not grid:
                continue

            # Encontrar el bloque de datos real (ignorar filas vacías iniciales)
            data_start = 0
            for ri, row in enumerate(grid):
                if any(str(c).strip() for c in row):
                    data_start = ri
                    break

            # Detectar fila de encabezado: primera fila con múltiples celdas llenas
            header_row = data_start
            for ri in range(data_start, min(data_start + 10, len(grid))):
                non_empty = sum(1 for c in grid[ri] if str(c).strip())
                if non_empty >= 2:
                    header_row = ri
                    break

            headers = [str(c or "").strip() for c in grid[header_row]]
            # Limpiar encabezados vacíos con índice
            headers = [h if h else f"Col{i+1}" for i, h in enumerate(headers)]

            data_rows = []
            for row in grid[header_row + 1:]:
                cleaned = [str(c or "").strip() for c in row]
                if any(c for c in cleaned):  # omitir filas completamente vacías
                    data_rows.append(cleaned)

            if data_rows:
                tables_extracted.append({
                    "source": f"excel_{sheet_name}",
                    "headers": headers,
                    "data": data_rows,
                    "sheet": sheet_name,
                    "page": None,
                })
                text_parts.append(_table_to_text(headers, data_rows, f"Hoja: {sheet_name}"))

        wb.close()
    except ImportError:
        logger.warning("openpyxl no disponible")
    except Exception as e:
        logger.warning(f"openpyxl error {path.name}: {e}")

    # ── pandas: fallback si openpyxl no leyó nada ───────────────────────
    if not tables_extracted:
        try:
            import pandas as pd
            xls = pd.ExcelFile(str(path))
            for sheet_name in xls.sheet_names:
                try:
                    # Intentar detectar la fila de encabezado
                    for header_row in [0, 1, 2, 3]:
                        try:
                            df = pd.read_excel(str(path), sheet_name=sheet_name,
                                               header=header_row, dtype=str)
                            df = df.dropna(how="all").fillna("")
                            if df.empty or len(df) < 1:
                                continue
                            headers = [str(c).strip() for c in df.columns]
                            rows = [list(r) for r in df.values.tolist()]
                            if any(any(c.strip() for c in r) for r in rows):
                                tables_extracted.append({
                                    "source": f"excel_{sheet_name}",
                                    "headers": headers,
                                    "data": rows,
                                    "sheet": sheet_name,
                                    "page": None,
                                })
                                text_parts.append(_table_to_text(headers, rows, f"Hoja: {sheet_name}"))
                                break
                        except Exception:
                            continue
                except Exception as se:
                    logger.debug(f"pandas hoja {sheet_name}: {se}")
        except ImportError:
            logger.warning("pandas no disponible")
        except Exception as e:
            logger.warning(f"pandas excel error: {e}")

    full_text = "\n\n".join(text_parts)
    result.update({
        "text": full_text,
        "tables": tables_extracted,
        "page_count": len(tables_extracted),
        "is_scanned": False,
        "pages": [{"page": 1, "text": full_text, "char_count": len(full_text)}],
        "error": None if tables_extracted else "No se pudo leer el Excel",
    })
    return result


def _unfuse_sheet(ws) -> list:
    """
    Convierte una hoja con celdas fusionadas en una grilla plana.
    Las celdas fusionadas propagan el valor de la celda superior-izquierda.
    """
    try:
        from openpyxl.utils import range_boundaries

        # Mapa de fusiones: (fila, col) → valor
        merged_values = {}
        for merged_range in ws.merged_cells.ranges:
            min_col, min_row, max_col, max_row = range_boundaries(str(merged_range))
            cell_value = ws.cell(row=min_row, column=min_col).value
            for r in range(min_row, max_row + 1):
                for c in range(min_col, max_col + 1):
                    merged_values[(r, c)] = cell_value

        rows = []
        for row in ws.iter_rows():
            row_data = []
            for cell in row:
                val = merged_values.get((cell.row, cell.column), cell.value)
                row_data.append(val if val is not None else "")
            rows.append(row_data)

        # Eliminar columnas completamente vacías
        if not rows:
            return rows
        n_cols = len(rows[0])
        non_empty_cols = [c for c in range(n_cols)
                          if any(str(rows[r][c]).strip() for r in range(len(rows)))]
        if not non_empty_cols:
            return []
        return [[row[c] for c in non_empty_cols] for row in rows]

    except Exception as e:
        logger.debug(f"_unfuse_sheet error: {e}")
        # Fallback: leer sin fusiones
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))
        return rows


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint (PPTX / PPT)
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pptx(path: Path) -> dict:
    result = _empty(file_type=path.suffix.lower())
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(path))
        text_parts = []
        tables_extracted = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                # Texto de formas
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # Tablas en diapositivas
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        if any(c for c in cells):
                            rows.append(cells)
                    if rows:
                        headers = rows[0] if rows else []
                        data = rows[1:] if len(rows) > 1 else []
                        tables_extracted.append({
                            "source": f"pptx_slide{slide_num}",
                            "headers": headers,
                            "data": data,
                            "page": slide_num,
                        })
                        slide_texts.append(_table_to_text(headers, data, f"Tabla slide {slide_num}"))

            if slide_texts:
                text_parts.append(f"[Diapositiva {slide_num}]\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(text_parts)
        result.update({
            "text": full_text,
            "tables": tables_extracted,
            "page_count": len(prs.slides),
            "is_scanned": False,
            "pages": [{"page": i+1, "text": t, "char_count": len(t)}
                      for i, t in enumerate(text_parts)],
            "error": None,
        })
    except ImportError:
        result["error"] = "python-pptx no disponible — instalar: pip install python-pptx"
    except Exception as e:
        result["error"] = str(e)
        logger.warning(f"PPTX parse error {path.name}: {e}")
    return result


# ─────────────────────────────────────────────────────────────────────────────
# Imágenes → marcado para OCR
# ─────────────────────────────────────────────────────────────────────────────

def _parse_image(path: Path) -> dict:
    result = _empty(file_type=path.suffix.lower())
    result.update({
        "is_scanned": True,   # → phase OCR lo procesará
        "page_count": 1,
        "error": None,
        "_needs_ocr": True,
    })
    return result


# ─────────────────────────────────────────────────────────────────────────────
# Utilidades
# ─────────────────────────────────────────────────────────────────────────────

def _table_to_text(headers: list, data: list, title: str = "") -> str:
    """
    Convierte una tabla en texto legible para el LLM.
    Formato: 'Encabezado: valor | Encabezado: valor' por fila.
    """
    lines = []
    if title:
        lines.append(f"\n--- {title} ---")

    clean_headers = [str(h).strip() for h in headers]

    for row in data:
        parts = []
        for ci, cell in enumerate(row):
            val = str(cell).strip()
            if not val or val.lower() in ("none", "nan", ""):
                continue
            hdr = clean_headers[ci] if ci < len(clean_headers) else f"Col{ci+1}"
            if hdr and hdr.lower() not in ("none", "nan", ""):
                parts.append(f"{hdr}: {val}")
            else:
                parts.append(val)
        if parts:
            lines.append(" | ".join(parts))

    return "\n".join(lines)


def _clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'\r\n|\r', '\n', text)
    # Eliminar líneas de solo espacios/guiones pero mantener contenido
    lines = []
    for line in text.split('\n'):
        stripped = line.strip()
        # Conservar si tiene al menos un carácter alfanumérico o símbolo útil
        if stripped and re.search(r'[\w$€%#@]', stripped):
            lines.append(stripped)
        elif stripped and len(stripped) <= 3:
            lines.append(stripped)
    return '\n'.join(lines)


def merge_document_texts(parsed_docs: list, max_chars: int = 150_000) -> str:
    """
    Combina textos de múltiples documentos.
    Orden de prioridad: estudio previo > invitación/pliego > presupuesto > cotización > otros.
    Inserta marcadores [DOC: archivo | PG: N] antes de cada página para que el LLM
    pueda citar el documento y número de página de cada dato extraído.
    """
    PRIORITY_KEYWORDS = {
        "estudio": 1, "previo": 1, "sector": 1,
        "invitacion": 2, "invitación": 2, "pliego": 2,
        "presupuesto": 3, "cotizacion": 3, "cotización": 3,
        "cdp": 4, "paa": 4,
    }

    def doc_priority(doc):
        name = (doc.get("filename") or "").lower()
        for kw, p in PRIORITY_KEYWORDS.items():
            if kw in name:
                return p
        return 5

    sorted_docs = sorted(parsed_docs, key=doc_priority)

    parts = []
    total = 0
    for doc in sorted_docs:
        fname = doc.get("filename", "documento")
        pages = doc.get("pages", [])

        # Build per-page chunks with source markers so the LLM can cite doc+page
        page_parts = []
        if pages:
            for pg in pages:
                page_num = pg.get("page", 1)
                page_text = pg.get("text", "").strip()
                if page_text:
                    page_parts.append(f"[DOC: {fname} | PG: {page_num}]\n{page_text}")
        else:
            # Fallback: use full text as a single page block
            text = doc.get("text", "").strip()
            if text:
                page_parts.append(f"[DOC: {fname} | PG: 1]\n{text}")

        if not page_parts:
            continue

        doc_header = f"\n{'='*70}\nDOCUMENTO: {fname}\n{'='*70}\n"
        chunk_body = "\n\n".join(page_parts)
        chunk = doc_header + chunk_body

        if total + len(chunk) > max_chars:
            remaining = max_chars - total - len(doc_header)
            if remaining > 500:
                parts.append(doc_header + chunk_body[:remaining] + "\n[... truncado ...]")
            break
        parts.append(chunk)
        total += len(chunk)

    return "\n".join(parts)


def extract_tables_summary(parsed_docs: list) -> list:
    """Retorna todas las tablas encontradas en todos los documentos."""
    all_tables = []
    for doc in parsed_docs:
        fname = doc.get("filename", "")
        for table in doc.get("tables", []):
            all_tables.append({**table, "filename": fname})
    return all_tables


def _empty(error=None, file_type="") -> dict:
    return {
        "text": "",
        "pages": [],
        "tables": [],
        "is_scanned": False,
        "page_count": 0,
        "file_type": file_type,
        "error": error,
    }
